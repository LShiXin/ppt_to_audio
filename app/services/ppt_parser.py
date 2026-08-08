import uuid
import shutil
from pathlib import Path
from pptx import Presentation

from app.config import UPLOADS_DIR


def parse_pptx(file_path: str) -> list[dict]:
    """
    Parse PPTX slides, extracting title, content, and notes.
    Uses python-pptx for accurate text extraction from shapes.
    Speaker notes are also extracted via python-pptx.
    """
    from pptx.shapes.placeholder import PP_PLACEHOLDER

    prs = Presentation(file_path)
    slides_data = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_number = slide_idx + 1

        title_texts = []
        content_texts = []
        notes_text = ""

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue

            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                if ph_type == PP_PLACEHOLDER.TITLE or ph_type == PP_PLACEHOLDER.CENTER_TITLE:
                    title_texts.append(text)
                elif ph_type == PP_PLACEHOLDER.SUBTITLE:
                    title_texts.append(text)
                else:
                    content_texts.append(text)
            else:
                content_texts.append(text)

        if title_texts:
            title = "\n".join(title_texts)
        elif content_texts:
            title = content_texts.pop(0)
            if len(title) > 100:
                content_texts.insert(0, title)
                title = ""
        else:
            title = ""

        content = "\n".join(content_texts)

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()

        slides_data.append({
            "slide_number": slide_number,
            "title": title,
            "content": content,
            "notes": notes_text,
        })

    return slides_data


def save_upload(file_content: bytes, original_filename: str, project_id: int) -> str:
    ext = Path(original_filename).suffix
    unique_name = f"{uuid.uuid4().hex}{ext}"
    project_dir = UPLOADS_DIR / str(project_id)
    project_dir.mkdir(parents=True, exist_ok=True)
    dest = project_dir / unique_name
    dest.write_bytes(file_content)
    return str(dest.relative_to(UPLOADS_DIR))


def extract_reference_text(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()
    if ext == ".txt" or ext == ".md":
        return Path(file_path).read_text(encoding="utf-8", errors="ignore")
    elif ext == ".pdf":
        return _extract_pdf_text(file_path)
    elif ext == ".docx":
        return _extract_docx_text(file_path)
    else:
        return ""


def _extract_pdf_text(file_path: str) -> str:
    try:
        import fitz
        doc = fitz.open(file_path)
        text = "\n".join(page.get_text() for page in doc)
        doc.close()
        return text
    except ImportError:
        return f"[PDF: {file_path}]"


def _extract_docx_text(file_path: str) -> str:
    try:
        from docx import Document
        doc = Document(file_path)
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except ImportError:
        return f"[DOCX: {file_path}]"


def extract_slide_images(pptx_path: str, output_dir: str) -> list[str]:
    """
    Render each slide as a PNG image using PyMuPDF for accurate font rendering.
    """
    import fitz

    out = Path(output_dir)
    out.mkdir(parents=True, exist_ok=True)

    doc = fitz.open(pptx_path)
    paths = []

    for i in range(len(doc)):
        page = doc[i]
        pix = page.get_pixmap(dpi=150)
        png_path = out / f"slide_{i + 1}.png"
        pix.save(str(png_path))
        paths.append(str(png_path.relative_to(out.parent)))

    doc.close()
    return paths
